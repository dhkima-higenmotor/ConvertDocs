import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def create_pptx_from_images(image_folder, output_pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    image_files = [f for f in os.listdir(image_folder) if f.lower().endswith('.png')]
    image_files.sort()

    for image_file in image_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_path = os.path.join(image_folder, image_file)
        
        with Image.open(img_path) as img:
            img_width, img_height = img.size

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Calculate the image size to fit the slide
        width_ratio = slide_width / img_width
        height_ratio = slide_height / img_height
        ratio = min(width_ratio, height_ratio)

        new_width = img_width * ratio
        new_height = img_height * ratio

        # Center the image
        left = (slide_width - new_width) / 2
        top = (slide_height - new_height) / 2

        slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

    prs.save(output_pptx_path)

def main():
    root = tk.Tk()
    root.withdraw()

    folder_path = filedialog.askdirectory(title="Select a directory containing PNG files")

    if folder_path:
        output_pptx_path = os.path.join(folder_path, "output.pptx")
        create_pptx_from_images(folder_path, output_pptx_path)
        print(f"PPTX file created successfully: {output_pptx_path}")

if __name__ == "__main__":
    main()
