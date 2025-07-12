import os
import tkinter as tk
from tkinter import filedialog, messagebox
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import shutil

def pdf_to_pptx():
    root = tk.Tk()
    root.withdraw()  # Hide the main window

    file_path = filedialog.askopenfilename(
        title="Select a PDF file",
        filetypes=[("PDF files", "*.pdf")]
    )

    if not file_path:
        return  # User cancelled the dialog

    pdf_dir = os.path.dirname(file_path)
    pdf_name_without_ext = os.path.splitext(os.path.basename(file_path))[0]
    
    # Create the temporary output subdirectory for PNGs
    temp_png_dir = os.path.join(pdf_dir, pdf_name_without_ext + "_temp_pngs")
    os.makedirs(temp_png_dir, exist_ok=True)

    png_files = []
    try:
        # --- Step 1: Convert PDF pages to PNG images ---
        pdf_document = fitz.open(file_path)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)  # Load page
            pix = page.get_pixmap()  # Render page to an image (pixmap)
            
            # Define output PNG file path
            output_png_path = os.path.join(temp_png_dir, f"{page_num + 1:03d}.png")
            
            # Save the pixmap as PNG
            pix.save(output_png_path)
            png_files.append(output_png_path)

        pdf_document.close()

        if not png_files:
            messagebox.showwarning("No Pages Found", "The selected PDF contains no pages to convert.")
            return

        # --- Step 2: Create PPTX from PNG images ---
        prs = Presentation()
        # Set slide size to a common aspect ratio (e.g., 16:9 widescreen)
        # Default is 10x7.5 inches (4:3). Let's use default and fit images.
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        # Sort PNG files to ensure correct page order
        png_files.sort()

        for png_path in png_files:
            slide_layout = prs.slide_layouts[6] # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Open image to get its dimensions
            img = Image.open(png_path)
            img_width_px, img_height_px = img.size
            img.close() # Close the image file

            # Calculate aspect ratios
            slide_aspect_ratio = slide_width_emu / slide_height_emu
            img_aspect_ratio = img_width_px / img_height_px

            # Determine new dimensions to fit image into slide while maintaining aspect ratio
            if img_aspect_ratio > slide_aspect_ratio:
                # Image is wider than slide, fit by width
                new_width_emu = slide_width_emu
                new_height_emu = int(slide_width_emu / img_aspect_ratio)
            else:
                # Image is taller than slide, fit by height
                new_height_emu = slide_height_emu
                new_width_emu = int(slide_height_emu * img_aspect_ratio)

            # Center the image on the slide
            left_emu = (slide_width_emu - new_width_emu) / 2
            top_emu = (slide_height_emu - new_height_emu) / 2

            # Add picture to slide
            slide.shapes.add_picture(str(png_path), left_emu, top_emu, width=new_width_emu, height=new_height_emu)

        output_pptx_path = os.path.join(pdf_dir, f"{pdf_name_without_ext}.pptx")
        prs.save(output_pptx_path)

        messagebox.showinfo("Success", f"PDF converted to PPTX successfully!\nSaved as: {output_pptx_path}")

    except fitz.FileNotFoundError:
        messagebox.showerror("Error", "Selected PDF file not found.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred during conversion: {e}")
    finally:
        # --- Step 3: Clean up temporary PNG files and directory ---
        if os.path.exists(temp_png_dir):
            try:
                shutil.rmtree(temp_png_dir)
            except Exception as e:
                messagebox.showwarning("Cleanup Warning", f"Could not remove temporary directory {temp_png_dir}: {e}")

if __name__ == "__main__":
    pdf_to_pptx()