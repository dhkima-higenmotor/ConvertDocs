import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from PIL import Image
import io

# Constants for EMU to Pixel conversion
EMU_PER_INCH = 914400
DEFAULT_DPI = 96

def emu_to_pixels(emu, dpi=DEFAULT_DPI):
    """Converts English Metric Units (EMU) to pixels."""
    return int(emu / EMU_PER_INCH * dpi)

def pptx_to_png():
    root = tk.Tk()
    root.withdraw()

    file_path = filedialog.askopenfilename(
        title="Select a PPTX file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )

    if not file_path:
        return

    try:
        presentation = Presentation(file_path)
        output_dir = os.path.join(os.path.dirname(file_path), os.path.splitext(os.path.basename(file_path))[0])
        os.makedirs(output_dir, exist_ok=True)

        # Get slide dimensions in pixels
        slide_width_px = emu_to_pixels(presentation.slide_width)
        slide_height_px = emu_to_pixels(presentation.slide_height)

        for i, slide in enumerate(presentation.slides):
            # Create a new transparent image to compose shapes on
            slide_image_canvas = Image.new('RGBA', (slide_width_px, slide_height_px), (0, 0, 0, 0))

            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        shape_img = Image.open(io.BytesIO(image_bytes))

                        # Convert shape position and size from EMU to pixels
                        left_px = emu_to_pixels(shape.left)
                        top_px = emu_to_pixels(shape.top)
                        width_px = emu_to_pixels(shape.width)
                        height_px = emu_to_pixels(shape.height)

                        if width_px > 0 and height_px > 0:
                            # Use a high-quality downsampling filter
                            shape_img = shape_img.resize((width_px, height_px), Image.Resampling.LANCZOS)
                            
                            # Ensure the shape image has an alpha channel for pasting
                            if shape_img.mode != 'RGBA':
                                shape_img = shape_img.convert('RGBA')
                            
                            # Paste the shape onto the canvas using its alpha channel
                            slide_image_canvas.paste(shape_img, (left_px, top_px), shape_img)

                    except Exception as e:
                        print(f"Could not process shape in slide {i+1}: {e}")

            # Create a final image with a white background
            final_image = Image.new("RGB", (slide_width_px, slide_height_px), "white")
            # Paste the composed slide canvas onto the white background
            # The mask is the alpha channel of the canvas
            final_image.paste(slide_image_canvas, (0, 0), slide_image_canvas)

            # Save the final image
            output_path = os.path.join(output_dir, f"{i+1:03d}.png")
            final_image.save(output_path, "PNG")

        messagebox.showinfo("Success", "All slides have been converted to PNG images.")

    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

if __name__ == "__main__":
    pptx_to_png()
